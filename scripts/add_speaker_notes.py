#!/usr/bin/env python3
"""Add speaker notes to 枢机验收答辩.pptx — one script per slide."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PPTX = REPO_ROOT / "assets" / "presentations" / "defense-2026" / "枢机验收答辩.pptx"

notes = [
    # ── Slide 1: Cover ──
    """各位老师好，我今天要汇报的项目是「枢机」—— 一个基于中国古代三省六部制的自动化软件开发系统。

我们把这个项目的核心思想概括为四句话：九卿协作、文档驱动、自我进化、全链路审计。

技术栈是 Rust + Tauri v2 做的桌面应用，目前完成了约四百个测试用例，处于原型阶段。""",
    # ── Slide 2: 架构总览 ──
    """我先介绍一下系统架构的核心设计。

首先是三省六部协作模型。我们设计了九个并发的 Actor，每个 Actor 用 tokio spawn 启动，通过 mpsc 信箱通信。三省是规划层——中书令做方案设计、门下侍中做审查、尚书令做执行调度。六部是执行层——吏、兵、工、刑、礼五个部门各司其职。另外还有两个子代理，一个做需求展开，一个做代码库调查。

第二个核心是文档驱动协作协议。我们定义了十几种文档类型，从需求到设计到计划到契约到报告，全部用 YAML 元数据加 Markdown 正文。部门之间不靠 LLM 对话上下文接力——所有的通信都通过读写这些结构化文档来完成。这个文档系统还有一个「朱批」机制——计划和审查类文档需要用户批准才能往下流转，三次不批就自动放行，避免系统卡死。

第三个是内阁的智能路由和自我进化。内阁根据用户输入的需求，自动判断复杂度，然后从十二个工作流技能里选一个最优的。它还有一个跨会话的记忆系统叫 Soul，分经验、教训、偏好三层，每次运行后会更新，下次启动时注入到 Prompt 里，相当于内阁越用越聪明。

第四个是审计追踪。所有的操作都有 JSONL 格式的审计日志，文件之间的依赖关系通过双向引用索引和谱系树来维护，任何修改都会自动保存 diff patch。如果修改了一个已经被下游引用的文档，系统会阻止——这叫不可变性检查。""",
    # ── Slide 3: 后端 Agent ──
    """这一页是组员 A 的工作——后端 Agent 实现与工具引擎。

上面这五张卡片是六部执行层的五个部门。吏部尚书负责详细设计，把中书令的方案拆成可执行的任务。兵部尚书负责写测试计划和接口契约。工部尚书负责 TDD 编码——它有一个很特别的「批次计划循环」，任务太大就拆成多个批次，一个批次一个批次地执行，规划的时候开推理，执行的时候关推理，通过 force_stop 来切换批次。刑部尚书负责运行测试，发现问题就提缺陷报告。礼部尚书负责规范检查和代码审计。

下面左边是工具注册与分发引擎。我们有九组工厂函数，根据角色的不同生成不同的工具定义——读文档的、读代码的、写文件的、操作文档的、审计的、执行 Shell 命令的，各不一样。所有工具的调用都经过 dispatch.rs 集中分发，在这里做审批状态检查、缓存失效，以及对命令和路径做安全过滤。

右边是 AgentController 控制循环。它驱动一次完整的 LLM 工具调用循环——step 一次就是一轮 API 交互。看门狗会检测异常行为：同一个工具重复调、只读不写、连续五次出错就自动停止。八个非内阁 Agent 共享同一个 runner.rs 执行框架，这个框架统一处理上下文压缩和检查点。""",
    # ── Slide 4: API 与基础设施 ──
    """组员 B 负责 API 客户端和上下文基础设施。

我们做了一个双格式的 API 客户端——同一个 AnthropicClient struct，如果 URL 里包含 anthropic.com，就走 Anthropic Messages API；否则走 OpenAI Chat Completions 格式。这意味着同一个代码可以接 DeepSeek、通义千问等各种模型。每个角色可以独立配置自己的 API Key、URL 和模型，还支持经济、均衡、质量三档预设。

中间是会话管理和上下文压缩。我们用了三层持久化存储——base prompt、soul prompt、context messages。当上下文超过阈值的时候，会把旧消息送给 LLM 生成摘要，保留最近二十四条。压缩策略分内阁专用和部门通用两种。而且支持在运行中途压缩——每二十次迭代触发一次，用原子写入来保证并发安全。

右边和下面两张卡片分别是检查点系统和安全沙箱。检查点用的是完全隔离的 git 仓库，跟项目的 .git 无关，每三百秒自动快照，恢复的时候 git stash 然后 checkout。Token 追踪是双轨的——一方面持久化 JSON 记录每次 API 调用的用量，另一方面在内存里实时追踪每轮会话的累计 Token 和缓存命中率。""",
    # ── Slide 5: 前端 ──
    """组员 C 负责前端驾驶舱界面的开发。

聊天系统是核心交互界面——ChatPanel 流式展示 Agent 输出，里面有个 ApprovalPromptCard 组件，把 LLM 输出的 option 标签解析成可点击的按钮，用户在聊天界面就能审批。DeptStatusPanel 实时显示九个 Agent 的状态，每个部门有不同颜色的光环动画，还有活动超时检测。

文档浏览器支持十几种文档类型的树形浏览，YAML 元数据直接在前端解析显示。审计面板可以看审计日志、谱系树、Diff 对比和违规报告。WorkflowTimeline 和 WorkflowGraph 用可视化方式展示工作流的步骤和部门之间的依赖关系。

右侧两块涵盖了面板系统和工程化。Token 用量可以按角色和时间窗口来查，检查点支持列表浏览和一键恢复。SettingsSidebar 统一管理 API 配置、上下文策略、Soul 编辑和外观主题。前端技术方面做了中英双语、三级字号调节、四套代码主题，十七个测试文件一百一十七个用例，用 Prettier 和 TypeScript 严格模式保证代码质量。""",
    # ── Slide 6: 测试 + 收尾 ──
    """最后一页是测试体系和程序演示。

测试方面，我们做到了百分之百的通过率。Rust 端九十六个单元测试加一百九十五个集成测试，十三个集成测试文件覆盖了二十二种路径攻击向量。前端十三套 Vitest 测试。还有三个端到端的 Mock LLM 场景。CI 方面有 Pre-commit Hook——提交前自动跑格式化、Clippy、单元测试和集成测试。所有测试都通过 TempDir 做文件隔离，用 test-threads 等于一避免状态竞争。

右边是程序演示。我们会启动 Tauri 桌面应用，现场输入一个需求，内阁自动判断复杂度和路由。然后三省六部依次协作执行，大家可以在部门状态面板上看到每个部门的实时状态变化。执行完成后在文档浏览器里查看产出——需求文档、设计文档、计划、契约代码、测试报告，全链路可追踪。

以上就是我们项目的全部汇报内容。谢谢大家，欢迎提问。""",
]


def main() -> None:
    pptx_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_PPTX
    if not pptx_path.is_file():
        print(f"找不到 PPTX 文件: {pptx_path}", file=sys.stderr)
        print("请将答辩 pptx 放到 assets/presentations/defense-2026/ 或传入路径参数。", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(pptx_path))

    for slide, note_text in zip(prs.slides, notes):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note_text

    prs.save(str(pptx_path))
    print(f"Done — added speaker notes to {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
